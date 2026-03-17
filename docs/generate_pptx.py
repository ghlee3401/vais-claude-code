#!/usr/bin/env python3
"""VAIS Code v0.11.0 프레젠테이션 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colors ──
BG       = RGBColor(0x0F, 0x17, 0x2A)
BG2      = RGBColor(0x1E, 0x29, 0x3B)
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2  = RGBColor(0x81, 0x8C, 0xF8)
ACCENT3  = RGBColor(0x34, 0xD3, 0x99)
WHITE    = RGBColor(0xF1, 0xF5, 0xF9)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
WARN     = RGBColor(0xFB, 0xBF, 0x24)
ERR      = RGBColor(0xF8, 0x71, 0x71)
BORDER   = RGBColor(0x33, 0x41, 0x55)
CARD_BG  = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def add_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_after=6, font_name='Arial'):
    p = tf.paragraphs[-1] if tf.paragraphs[0].text == '' and len(tf.paragraphs) == 1 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    return p

def new_para(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_after=6, font_name='Arial'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    return p

def add_card(slide, left, top, width, height, accent_color=ACCENT):
    """Add a rounded card with accent left border"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.05), Inches(0.06), Inches(height - 0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    return shape

def add_table(slide, left, top, width, rows, cols, col_widths=None):
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(rows * 0.45))
    tbl = tbl_shape.table
    for i, row in enumerate(tbl.rows):
        for j, cell in enumerate(row.cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG2 if i == 0 else BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = ACCENT if i == 0 else WHITE
                p.font.bold = (i == 0)
                p.font.name = 'Arial'
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(w)
    return tbl

def set_cell(tbl, row, col, text, color=None, size=13, bold=False):
    cell = tbl.cell(row, col)
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.color.rgb = color or (ACCENT if row == 0 else WHITE)
        p.font.bold = bold or (row == 0)
        p.font.name = 'Arial'

# ════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

tb = add_textbox(slide, 1.2, 1.5, 11, 1.5)
add_text(tb.text_frame, 'VAIS Code', size=72, color=ACCENT, bold=True, alignment=PP_ALIGN.LEFT)
tb2 = add_textbox(slide, 1.2, 3.2, 11, 1)
add_text(tb2.text_frame, 'v0.11.0 기술 소개', size=44, color=WHITE, bold=True)
tb3 = add_textbox(slide, 1.2, 4.5, 9, 0.8)
add_text(tb3.text_frame, 'Claude Code 플러그인 기반 9단계 AI 개발 워크플로우', size=24, color=MUTED)
tb4 = add_textbox(slide, 1.2, 5.5, 10, 0.5)
add_text(tb4.text_frame, '체이닝 문법 · 병렬 에이전트 · 피드백 루프 · Gap 분석 · Manager Memory', size=16, color=MUTED)
tb5 = add_textbox(slide, 1.2, 6.3, 3, 0.4)
add_text(tb5.text_frame, '2026.03', size=14, color=MUTED)

# ════════════════════════════════════════════
# SLIDE 2: 목차
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.4, 6, 0.8)
add_text(tb.text_frame, '📑  목차', size=36, color=ACCENT, bold=True)

sections = [
    ('Part 1 — 개요', ['VAIS Code란?', '핵심 가치와 차별점', '전체 아키텍처'], ACCENT),
    ('Part 2 — 워크플로우', ['9단계 파이프라인', '체이닝 문법 (: 순차, + 병렬)', '게이트 & 핑퐁 피드백 루프'], ACCENT2),
    ('Part 3 — 시스템 구조', ['에이전트 팀 구성', '훅(Hook) 시스템', '상태 관리 & Memory'], ACCENT3),
    ('Part 4 — 실전 활용', ['파일 구조 & 데모 흐름', 'Gap 분석 & /vais fix', 'v0.11.0 변경 사항'], WARN),
]

for i, (title, items, color) in enumerate(sections):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.6 + row * 2.8
    card = add_card(slide, x, y, 5.8, 2.4, color)
    tb = add_textbox(slide, x + 0.3, y + 0.2, 5.2, 0.5)
    add_text(tb.text_frame, title, size=22, color=color, bold=True)
    for item in items:
        tb_item = add_textbox(slide, x + 0.5, y + 0.7 + items.index(item) * 0.45, 5, 0.4)
        add_text(tb_item.text_frame, f'→  {item}', size=16, color=WHITE)

# ════════════════════════════════════════════
# SLIDE 3: Section - 개요
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 2, 1.5, 9, 2)
add_text(tb.text_frame, '01', size=120, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(slide, 2, 4.0, 9, 1)
add_text(tb2.text_frame, '개요', size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
tb3 = add_textbox(slide, 2, 5.2, 9, 0.6)
add_text(tb3.text_frame, 'VAIS Code는 무엇이고, 왜 필요한가?', size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 4: What is VAIS Code
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.4, 10, 0.8)
add_text(tb.text_frame, 'VAIS Code란?', size=36, color=ACCENT, bold=True)
tb2 = add_textbox(slide, 0.8, 1.2, 11, 0.7)
add_text(tb2.text_frame, 'Claude Code 플러그인으로 동작하는 9단계 AI 개발 워크플로우 도구', size=20, color=WHITE)

cards_data = [
    ('🎯 문제', 'AI 코딩 도구에게 "만들어줘" 하면\n기획 없이 바로 코딩\n→ 요구사항 누락, 설계 없는 코드', ACCENT),
    ('💡 해결', '조사→기획→설계→구현→검증까지\n단계별 산출물을 강제하여\n체계적 개발 워크플로우 보장', ACCENT2),
    ('✨ 결과', '기획서 기반 코딩 규칙 준수\n설계 대비 Gap 분석으로\n90%+ 일치율 자동 검증', ACCENT3),
]
for i, (title, desc, color) in enumerate(cards_data):
    x = 0.8 + i * 4.1
    card = add_card(slide, x, 2.1, 3.8, 3.0, color)
    tb = add_textbox(slide, x + 0.3, 2.3, 3.3, 0.5)
    add_text(tb.text_frame, title, size=22, color=color, bold=True)
    tb2 = add_textbox(slide, x + 0.3, 2.9, 3.3, 2.0)
    add_text(tb2.text_frame, desc, size=15, color=WHITE, space_after=4)

badges = add_textbox(slide, 0.8, 5.5, 12, 0.5)
add_text(badges.text_frame, 'Claude Code Plugin  ·  MIT License  ·  한국어 네이티브  ·  v0.11.0', size=14, color=MUTED)

# ════════════════════════════════════════════
# SLIDE 5: 핵심 가치
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.4, 10, 0.8)
add_text(tb.text_frame, '핵심 가치 & 차별점', size=36, color=ACCENT, bold=True)

values = [
    ('📋 기획 없이 코드 금지', '기획서가 없으면 코딩 시작 불가.\n기획서에 코딩 규칙, 기능 목록, 정책이\n포함되어 SSOT(단일 진실 소스) 역할', ACCENT),
    ('🔗 체이닝 문법', 'plan:ia:wireframe → 순차 실행\nfe+be → 병렬 실행\n직관적 DSL로 워크플로우 조합', ACCENT2),
    ('🔎 Gap 분석 자동 검증', '설계 vs 구현 자동 비교\n일치율 90% 미만 시 최대 5회 반복\nOWASP Top 10 보안 검사 포함', ACCENT3),
    ('🧠 Manager Memory', '모든 의사결정, 피처 간 의존성,\n기술 부채를 영속 저장.\n다음 세션에서도 맥락 유지', WARN),
]
for i, (title, desc, color) in enumerate(values):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.5 + row * 2.8
    card = add_card(slide, x, y, 5.8, 2.4, color)
    tb = add_textbox(slide, x + 0.3, y + 0.2, 5.2, 0.5)
    add_text(tb.text_frame, title, size=20, color=color, bold=True)
    tb2 = add_textbox(slide, x + 0.3, y + 0.8, 5.2, 1.4)
    add_text(tb2.text_frame, desc, size=15, color=WHITE, space_after=4)

# ════════════════════════════════════════════
# SLIDE 6: Architecture
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.4, 10, 0.8)
add_text(tb.text_frame, '전체 아키텍처', size=36, color=ACCENT, bold=True)

diagram_text = (
    '사용자\n'
    '    │  /vais auto login\n'
    '    ▼\n'
    '┌──────────────┐\n'
    '│  SKILL.md     │ ← 스킬 진입점 (액션 라우팅)\n'
    '└──────┬───────┘\n'
    '       ▼\n'
    '┌──────────┐  ┌──────────┐  ┌─────────────┐\n'
    '│ Manager  │→ │ Tech Lead│→ │ 에이전트 팀  │\n'
    '│ (opus)   │  │ (opus)   │  │ (sonnet)    │\n'
    '│ memory   │  │ 오케스트 │  │ designer    │\n'
    '└──────────┘  └──────────┘  │ fe/be dev   │\n'
    '      ▲            │       │ reviewer    │\n'
    '      │        ┌───┴───┐  └─────────────┘\n'
    '  피드백/기록  │ Hooks │  │ Status │\n'
    '              │ 5개 훅│  │ .vais/ │\n'
    '              └───────┘  └────────┘'
)
box = add_card(slide, 1.5, 1.5, 10, 5.2, ACCENT)
tb2 = add_textbox(slide, 1.8, 1.7, 9.5, 5.0)
add_text(tb2.text_frame, diagram_text, size=14, color=WHITE, font_name='Courier New', space_after=2)

# ════════════════════════════════════════════
# SLIDE 7: Section - 워크플로우
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 2, 1.5, 9, 2)
add_text(tb.text_frame, '02', size=120, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(slide, 2, 4.0, 9, 1)
add_text(tb2.text_frame, '워크플로우', size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
tb3 = add_textbox(slide, 2, 5.2, 9, 0.6)
add_text(tb3.text_frame, '9단계 파이프라인, 체이닝, 게이트', size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 8: 9단계 파이프라인
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '9단계 개발 파이프라인', size=32, color=ACCENT, bold=True)

# Flow
flow = '🔭 조사·탐색 → 📋 기획 → 🗺 IA → 🖼 와이어프레임 → 🎨 설계 → 💻 FE → ⚙️ BE → 🔎 Gap → 🔍 검토'
tb2 = add_textbox(slide, 0.8, 1.1, 12, 0.5)
add_text(tb2.text_frame, flow, size=16, color=ACCENT3, bold=True)

# Table
tbl = add_table(slide, 0.8, 1.8, 11.7, 10, 4, [1.8, 3.5, 3.0, 3.4])
headers = ['단계', '산출물', '에이전트', '비고']
for j, h in enumerate(headers):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)

rows_data = [
    ('🔭 research', 'docs/01-research/', 'Tech Lead 직접', '아이디어 → 기능 도출'),
    ('📋 plan', 'docs/02-plan/ + .vais/features/', 'Tech Lead 직접', '⚡ 게이트 | 코딩 규칙 포함'),
    ('🗺 ia', 'docs/03-ia/', 'designer', '사이트맵, 네비게이션'),
    ('🖼 wireframe', 'docs/04-wireframe/', 'designer', 'ASCII + HTML 와이어프레임'),
    ('🎨 design', 'docs/05-design/', 'designer + backend-dev', '🔀 병렬 | UI + DB 설계'),
    ('💻 fe', 'src/', 'frontend-dev', '⚡ 게이트 | 🔀 병렬'),
    ('⚙️ be', 'src/ or server/', 'backend-dev', '🔀 병렬 (fe와 동시)'),
    ('🔎 check', 'docs/06-check/', 'reviewer', 'Gap 분석 + 보안 검사'),
    ('🔍 review', 'docs/07-review/', 'reviewer', 'QA + 최종 리뷰'),
]
for i, (a, b, c, d) in enumerate(rows_data):
    set_cell(tbl, i+1, 0, a)
    set_cell(tbl, i+1, 1, b)
    set_cell(tbl, i+1, 2, c)
    set_cell(tbl, i+1, 3, d)

# ════════════════════════════════════════════
# SLIDE 9: 체이닝 문법
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '체이닝 문법', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, ': (순차) + (병렬)을 조합하여 워크플로우 커스터마이징', size=18, color=MUTED)

# Left: grammar table
tbl = add_table(slide, 0.8, 1.6, 5.5, 4, 3, [1.0, 1.5, 3.0])
for j, h in enumerate(['기호', '의미', '예시']):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)
set_cell(tbl, 1, 0, ':');  set_cell(tbl, 1, 1, '순차 실행');  set_cell(tbl, 1, 2, 'plan:ia:wireframe')
set_cell(tbl, 2, 0, '+');  set_cell(tbl, 2, 1, '병렬 실행');  set_cell(tbl, 2, 2, 'fe+be')
set_cell(tbl, 3, 0, '혼합'); set_cell(tbl, 3, 1, '순차+병렬'); set_cell(tbl, 3, 2, 'plan:design:fe+be:check')

# Right: examples
examples = (
    '# 단일 실행\n'
    '/vais plan 로그인기능\n\n'
    '# 순차 체이닝: 기획→IA→와이어프레임\n'
    '/vais plan:ia:wireframe 로그인기능\n\n'
    '# 병렬: 프론트 + 백엔드 동시\n'
    '/vais fe+be 로그인기능\n\n'
    '# 혼합: 기획→설계→구현(병렬)→검증\n'
    '/vais plan:design:fe+be:check 로그인기능\n\n'
    '# 전체 자동\n'
    '/vais auto 로그인기능\n\n'
    '# 범위 지정 (한글)\n'
    '"기획부터 백엔드까지"\n'
    '→ plan:ia:wireframe:design:fe+be'
)
card = add_card(slide, 6.8, 1.6, 5.8, 5.3, ACCENT2)
tb3 = add_textbox(slide, 7.1, 1.8, 5.3, 5.0)
add_text(tb3.text_frame, examples, size=13, color=WHITE, font_name='Courier New', space_after=2)

# Level table
tbl2 = add_table(slide, 0.8, 3.6, 5.5, 3, 3, [1.2, 2.5, 1.8])
for j, h in enumerate(['레벨', '범위', '용도']):
    set_cell(tbl2, 0, j, h, ACCENT, bold=True)
set_cell(tbl2, 1, 0, 'Quick'); set_cell(tbl2, 1, 1, 'plan → check (7단계)'); set_cell(tbl2, 1, 2, 'MVP, 프로토타입')
set_cell(tbl2, 2, 0, 'Full');  set_cell(tbl2, 2, 1, 'research → review (9)'); set_cell(tbl2, 2, 2, '프로덕션')

# ════════════════════════════════════════════
# SLIDE 10: 게이트 & 핑퐁
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '게이트 & 핑퐁 피드백 루프', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, 'v0.11.0 핵심 변경 — 사용자가 "계속"을 누를 때까지 멈춤', size=18, color=MUTED)

# Left: gate info
card_l = add_card(slide, 0.8, 1.6, 5.5, 2.2, WARN)
tb2 = add_textbox(slide, 1.1, 1.8, 5.0, 0.5)
add_text(tb2.text_frame, '게이트 체크포인트', size=20, color=WARN, bold=True)
tb3 = add_textbox(slide, 1.1, 2.4, 5.0, 1.2)
add_text(tb3.text_frame, 'vais.config.json의 orchestration.gates에서 설정\n기본값: plan, fe\n\n• plan 게이트 — 기획 완료 후 방향 확인\n• fe 게이트 — 구현 시작 전 설계 결과 확인', size=14, color=WHITE, space_after=3)

# Right: ping-pong diagram
pingpong = (
    '게이트 도달\n'
    '  │\n'
    '  ▼\n'
    'AskUserQuestion\n'
    '  ├─ "계속"      → 다음 단계 ✅\n'
    '  ├─ "수정 요청" → 피드백 수신\n'
    '  │     │\n'
    '  │     ▼\n'
    '  │   피드백 반영 → 문서/코드 수정\n'
    '  │     │\n'
    '  │     ▼\n'
    '  │   📝 수정 결과 요약 출력\n'
    '  │     │\n'
    '  │     ▼\n'
    '  │   AskUserQuestion (재확인)\n'
    '  │     ├─ "계속"      → 다음 단계 ✅\n'
    '  │     ├─ "추가 수정" → ↑ 반복 🔄\n'
    '  │     └─ "중단"      → 저장 후 종료 🛑\n'
    '  │\n'
    '  └─ "중단"      → 저장 후 종료 🛑'
)
card_r = add_card(slide, 6.8, 1.6, 5.8, 5.3, ACCENT3)
tb4 = add_textbox(slide, 7.1, 1.7, 5.3, 0.5)
add_text(tb4.text_frame, '핑퐁 루프 (v0.11.0 신규)', size=20, color=ACCENT3, bold=True)
tb5 = add_textbox(slide, 7.1, 2.3, 5.3, 4.5)
add_text(tb5.text_frame, pingpong, size=12, color=WHITE, font_name='Courier New', space_after=1)

# Key point
card_b = add_card(slide, 0.8, 4.2, 5.5, 1.5, ERR)
tb6 = add_textbox(slide, 1.1, 4.4, 5.0, 1.2)
add_text(tb6.text_frame, '⚠️  중요\n사용자가 명시적으로 "계속"을 선택할 때까지\n절대 다음 단계로 넘어가지 않습니다.', size=15, color=WHITE, bold=False, space_after=3)

# ════════════════════════════════════════════
# SLIDE 11: Section - 시스템 구조
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 2, 1.5, 9, 2)
add_text(tb.text_frame, '03', size=120, color=ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(slide, 2, 4.0, 9, 1)
add_text(tb2.text_frame, '시스템 구조', size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
tb3 = add_textbox(slide, 2, 5.2, 9, 0.6)
add_text(tb3.text_frame, '에이전트, 훅, 상태 관리', size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 12: Agent Team
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '에이전트 팀 구성', size=32, color=ACCENT, bold=True)

agents = [
    ('📋 Manager', 'opus\n프로젝트 기억\n전략적 판단\n크로스-피처 분석\nTech Lead 지시', WARN),
    ('🏗 Tech Lead', 'opus\n아키텍처 결정\n팀 오케스트레이션\n게이트 관리\n품질 관리', ACCENT),
    ('🎨 Designer', 'sonnet\nIA 설계\n와이어프레임\nUI/UX 설계\n디자인 토큰', ACCENT2),
    ('💻 FE / ⚙️ BE', 'sonnet\n프론트엔드 구현\n백엔드 API\nDB 설계/연동\n테스트 코드', ACCENT3),
    ('🔍 Reviewer', 'sonnet\nGap 분석\nOWASP 보안\nQA 검증\n코드 리뷰', ERR),
]
for i, (name, desc, color) in enumerate(agents):
    x = 0.5 + i * 2.5
    card = add_card(slide, x, 1.2, 2.3, 3.5, color)
    tb2 = add_textbox(slide, x + 0.2, 1.4, 2.0, 0.5)
    add_text(tb2.text_frame, name, size=15, color=color, bold=True)
    tb3 = add_textbox(slide, x + 0.2, 1.9, 2.0, 2.5)
    add_text(tb3.text_frame, desc, size=12, color=WHITE, space_after=2)

# Orchestration table
tbl = add_table(slide, 0.8, 5.0, 11.7, 5, 3, [2.0, 5.0, 4.7])
for j, h in enumerate(['패턴', '설명', '적용 단계']):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)
patterns = [
    ('leader', 'Tech Lead가 직접 수행', 'research, plan'),
    ('delegate', '전문 에이전트에 위임', 'ia, wireframe, fe, be'),
    ('parallel', '여러 에이전트 동시 호출', 'design(UI+DB), fe+be'),
    ('council', 'reviewer + 리드 공동 검토', 'check, review'),
]
for i, (a, b, c) in enumerate(patterns):
    set_cell(tbl, i+1, 0, a, ACCENT3)
    set_cell(tbl, i+1, 1, b)
    set_cell(tbl, i+1, 2, c)

# ════════════════════════════════════════════
# SLIDE 13: 비용 최적화
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '모델 비용 최적화', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, '판단은 Opus, 구현은 Sonnet — 역할별 모델 분리', size=18, color=MUTED)

model_diagram = (
    'opus  ━━━━━━━━━━━━━━━━━━━━\n'
    '  │  판단, 전략, 아키텍처\n'
    '  ├── Manager    (프로젝트 기억/판단)\n'
    '  └── Tech Lead  (아키텍처/오케스트레이션)\n'
    '\n'
    'sonnet ━━━━━━━━━━━━━━━━━━━━\n'
    '  │  코드 구현, 문서 작성\n'
    '  ├── Designer      (IA/와이어프레임/설계)\n'
    '  ├── Frontend Dev  (FE 구현)\n'
    '  ├── Backend Dev   (BE 구현 + DB)\n'
    '  └── Reviewer      (Gap/QA/리뷰)'
)
card = add_card(slide, 0.8, 1.6, 6.0, 4.5, ACCENT)
tb2 = add_textbox(slide, 1.1, 1.8, 5.5, 4.2)
add_text(tb2.text_frame, model_diagram, size=15, color=WHITE, font_name='Courier New', space_after=2)

card2 = add_card(slide, 7.3, 1.6, 5.3, 2.0, ACCENT3)
tb3 = add_textbox(slide, 7.6, 1.8, 4.8, 0.5)
add_text(tb3.text_frame, '핵심 원리', size=20, color=ACCENT3, bold=True)
tb4 = add_textbox(slide, 7.6, 2.4, 4.8, 1.0)
add_text(tb4.text_frame, '→ 기획/아키텍처: 정확성 우선 → opus\n→ 코드 구현: 속도/비용 효율 → sonnet\n→ 단순 탐색: haiku 활용 가능', size=15, color=WHITE, space_after=3)

card3 = add_card(slide, 7.3, 4.0, 5.3, 2.0, WARN)
tb5 = add_textbox(slide, 7.6, 4.2, 4.8, 0.5)
add_text(tb5.text_frame, '비용 절감 효과', size=20, color=WARN, bold=True)
tb6 = add_textbox(slide, 7.6, 4.8, 4.8, 1.0)
add_text(tb6.text_frame, '9단계 중 2단계만 opus 사용.\n나머지 7단계는 sonnet으로 처리하여\n전체 비용 60~70% 절감 (추정)', size=15, color=WHITE, space_after=3)

# ════════════════════════════════════════════
# SLIDE 14: Hook System
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '훅(Hook) 시스템', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, 'Claude Code 이벤트에 자동 반응하는 5개 훅', size=18, color=MUTED)

tbl = add_table(slide, 0.8, 1.6, 11.7, 6, 4, [2.2, 2.0, 3.5, 4.0])
for j, h in enumerate(['훅', '이벤트', '스크립트', '역할']):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)

hooks_data = [
    ('SessionStart', '세션 시작', 'session-start.js', '.vais/ 생성, 피처 상태 표시'),
    ('PreToolUse', 'Bash 실행 전', 'bash-guard.js', 'rm -rf, DROP, sudo 등 차단'),
    ('PostToolUse', 'Write|Edit 후', 'doc-tracker.js', '문서 감지 → 상태 자동 업데이트'),
    ('UserPromptSubmit', '사용자 입력', 'prompt-handler.js', '의도 감지 → /vais 명령 제안'),
    ('Stop', '응답 완료', 'stop-handler.js', '진행률 바 + 다음 단계 안내'),
]
for i, (a, b, c, d) in enumerate(hooks_data):
    set_cell(tbl, i+1, 0, a, ACCENT2)
    set_cell(tbl, i+1, 1, b)
    set_cell(tbl, i+1, 2, c)
    set_cell(tbl, i+1, 3, d)

card = add_card(slide, 0.8, 4.6, 11.7, 1.2, ACCENT3)
tb2 = add_textbox(slide, 1.1, 4.8, 11.0, 0.9)
add_text(tb2.text_frame, 'v0.11.0 변경: no-op 훅 2개 (notification.js, subagent-stop.js) 제거 → 불필요한 오버헤드 제거', size=15, color=WHITE)

# ════════════════════════════════════════════
# SLIDE 15: 상태 관리
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '상태 관리 시스템', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, '.vais/ 디렉토리에 모든 상태를 JSON으로 영속 관리', size=18, color=MUTED)

file_tree = (
    '.vais/\n'
    '├── status.json        # 워크플로우 상태\n'
    '│   ├── activeFeature\n'
    '│   └── features\n'
    '│       └── login\n'
    '│           ├── currentPhase\n'
    '│           ├── phases (9개)\n'
    '│           └── gapAnalysis\n'
    '│\n'
    '├── memory.json        # Manager Memory\n'
    '│   └── entries[]\n'
    '│       ├── decision / change\n'
    '│       ├── feedback / dependency\n'
    '│       ├── debt / error\n'
    '│       └── milestone\n'
    '│\n'
    '├── features/          # 피처 레지스트리\n'
    '│   └── login.json\n'
    '│\n'
    '└── project-config.json'
)
card = add_card(slide, 0.8, 1.6, 6.0, 5.2, ACCENT)
tb2 = add_textbox(slide, 1.1, 1.8, 5.5, 5.0)
add_text(tb2.text_frame, file_tree, size=13, color=WHITE, font_name='Courier New', space_after=1)

# Module table
tbl = add_table(slide, 7.3, 1.6, 5.3, 6, 2, [1.5, 3.8])
for j, h in enumerate(['모듈', '역할']):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)
mods = [
    ('paths.js', '경로 레지스트리, loadConfig'),
    ('status.js', '워크플로우 상태, updatePhase'),
    ('memory.js', 'Manager Memory, addEntry'),
    ('io.js', '훅 I/O, outputAllow/Block'),
    ('debug.js', '디버그 로깅 (VAIS_DEBUG=1)'),
]
for i, (a, b) in enumerate(mods):
    set_cell(tbl, i+1, 0, a, ACCENT3)
    set_cell(tbl, i+1, 1, b)

# ════════════════════════════════════════════
# SLIDE 16: Memory System
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, 'Manager Memory 시스템', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, '프로젝트의 모든 의사결정과 컨텍스트를 영속 기억', size=18, color=MUTED)

tbl = add_table(slide, 0.8, 1.6, 6.5, 8, 3, [1.5, 2.0, 3.0])
for j, h in enumerate(['타입', '용도', '예시']):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)
entries = [
    ('decision', '기술/비즈니스 결정', 'React 대신 Next.js — SSR'),
    ('change', 'fix로 인한 변경', 'JWT → 세션 변경'),
    ('feedback', '게이트 사용자 피드백', '소셜 로그인도 추가'),
    ('dependency', '피처 간 의존', 'cart → login 의존'),
    ('debt', '기술 부채', '에러 핸들링 미완'),
    ('error', '실패/재시도', '빌드 실패 — 타입 에러'),
    ('milestone', '단계 완료', 'plan 단계 완료'),
]
for i, (a, b, c) in enumerate(entries):
    set_cell(tbl, i+1, 0, a, ACCENT2)
    set_cell(tbl, i+1, 1, b)
    set_cell(tbl, i+1, 2, c)

scenarios = [
    ('🔍 Query 모드', '"프로젝트 현황 알려줘"\n→ memory 조회 → 피처 목록,\n진행률, 최근 활동 요약', ACCENT),
    ('⚡ 크로스-피처 분석', '"login 수정해줘"\n→ dependency 맵 조회\n→ cart, payment 영향 알림', ACCENT2),
    ('📋 기술 부채 관리', '미해결 부채 목록 조회\n해결 시 resolveDebt()로\n표시', ACCENT3),
]
for i, (title, desc, color) in enumerate(scenarios):
    y = 1.6 + i * 1.8
    card = add_card(slide, 7.8, y, 4.8, 1.5, color)
    tb2 = add_textbox(slide, 8.1, y + 0.15, 4.3, 0.4)
    add_text(tb2.text_frame, title, size=15, color=color, bold=True)
    tb3 = add_textbox(slide, 8.1, y + 0.55, 4.3, 0.8)
    add_text(tb3.text_frame, desc, size=12, color=WHITE, space_after=2)

# ════════════════════════════════════════════
# SLIDE 17: Section - 실전 활용
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 2, 1.5, 9, 2)
add_text(tb.text_frame, '04', size=120, color=WARN, bold=True, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(slide, 2, 4.0, 9, 1)
add_text(tb2.text_frame, '실전 활용', size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
tb3 = add_textbox(slide, 2, 5.2, 9, 0.6)
add_text(tb3.text_frame, '파일 구조, 데모 흐름, 변경 사항', size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 18: Demo Flow
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '실행 데모 — /vais auto login', size=32, color=ACCENT, bold=True)

demo = (
    '1. Manager — memory.json 조회 → 기존 피처 의존성 확인\n\n'
    '2. Tech Lead — research → docs/01-research/login.md\n\n'
    '3. Tech Lead — plan  ⚡게이트\n'
    '   기획서 → 피처 레지스트리 → 핑퐁 루프\n'
    '   └─ 사용자: "소셜 로그인도 넣어줘" → 반영 → "계속"\n\n'
    '4. designer — ia → 사이트맵 + 네비게이션\n\n'
    '5. designer — wireframe → ASCII/HTML 와이어프레임\n\n'
    '6. 🔀 병렬 — design\n'
    '   ├─ designer → UI 설계 (디자인 토큰, 컴포넌트)\n'
    '   └─ backend-dev → DB 설계 (ERD, 스키마)\n\n'
    '7. 🔀 병렬 — fe+be  ⚡게이트 (fe 시작 전)\n'
    '   ├─ frontend-dev → React 컴포넌트 구현\n'
    '   └─ backend-dev → API + DB 연동\n\n'
    '8. reviewer — check → Gap 분석 (90% 미만 시 자동 반복)\n\n'
    '9. reviewer — review → QA + OWASP 보안 + 최종 리뷰'
)
card = add_card(slide, 0.8, 1.2, 11.7, 5.8, ACCENT)
tb2 = add_textbox(slide, 1.1, 1.3, 11.2, 5.6)
add_text(tb2.text_frame, demo, size=14, color=WHITE, font_name='Courier New', space_after=1)

# ════════════════════════════════════════════
# SLIDE 19: Gap Analysis
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, 'Gap 분석 — 자동 검증 루프', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, '설계 문서 vs 구현 코드를 자동 비교, 90% 미만이면 자동 재시도', size=18, color=MUTED)

# Left: process
process = (
    '1. 피처 레지스트리에서 요구사항 추출\n'
    '2. Glob/Grep으로 구현 코드 탐색\n'
    '3. 각 요구사항의 구현 여부 판정\n'
    '4. 일치율 산출: (구현 수/전체) × 100\n'
    '5. OWASP Top 10 보안 검사\n'
    '6. QA 시나리오 생성\n'
    '7. Gap 방향 판단: 코드 수정? 스펙 수정?\n'
    '8. 90% 미만 → 최대 5회 자동 반복'
)
card = add_card(slide, 0.8, 1.6, 5.5, 3.5, ACCENT3)
tb2 = add_textbox(slide, 1.1, 1.8, 5.0, 3.2)
add_text(tb2.text_frame, process, size=15, color=WHITE, space_after=4)

# Right: diagram
gap_diagram = (
    'check 시작\n'
    '  │\n'
    '  ▼\n'
    '빌드 검증 (npm build / tsc)\n'
    '  │\n'
    '  ▼\n'
    'Gap 분석 실행\n'
    '  │\n'
    '  ├─ ≥ 90% → 통과 ✅\n'
    '  │\n'
    '  └─ < 90%\n'
    '       │\n'
    '       ▼\n'
    '     Gap 방향 판단\n'
    '     ├─ 코드 누락 → 코드 수정\n'
    '     └─ 스펙 과잉 → 사용자 확인\n'
    '       │\n'
    '     반복 (iteration++)\n'
    '     ├─ ≤ 5 → 재분석 🔄\n'
    '     └─ > 5 → 중단 🛑'
)
card2 = add_card(slide, 6.8, 1.6, 5.8, 5.3, ACCENT)
tb3 = add_textbox(slide, 7.1, 1.8, 5.3, 5.0)
add_text(tb3.text_frame, gap_diagram, size=13, color=WHITE, font_name='Courier New', space_after=1)

# Config
card3 = add_card(slide, 0.8, 5.4, 5.5, 1.4, WARN)
tb4 = add_textbox(slide, 1.1, 5.5, 5.0, 1.2)
add_text(tb4.text_frame, '설정 (vais.config.json)\nmatchThreshold: 90  |  maxIterations: 5  |  autoIterate: true', size=13, color=WHITE, space_after=2)

# ════════════════════════════════════════════
# SLIDE 20: /vais fix
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '/vais fix — 영향 분석 기반 수정', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, 'Manager가 영향 범위를 분석 후 적절한 단계 체이닝', size=18, color=MUTED)

fix_flow = (
    '사용자: "로그인 화면 디자인 바꿔줘"\n'
    '  │\n'
    '  ▼\n'
    'Manager\n'
    '  ├── memory 조회 (과거 결정)\n'
    '  ├── dependency 맵 확인\n'
    '  └── 영향 범위 분석\n'
    '  │\n'
    '  ▼\n'
    '영향 분류\n'
    '  ├── UI만   → wireframe:design:fe\n'
    '  ├── API만  → design:be\n'
    '  ├── 전체   → design:fe+be:check\n'
    '  └── 정책   → plan:design:fe+be:check\n'
    '  │\n'
    '  ▼\n'
    'Tech Lead → 팀 실행\n'
    '  │\n'
    '  ▼\n'
    'Manager → memory 기록 (change)'
)
card = add_card(slide, 0.8, 1.6, 6.0, 5.2, ACCENT)
tb2 = add_textbox(slide, 1.1, 1.8, 5.5, 5.0)
add_text(tb2.text_frame, fix_flow, size=14, color=WHITE, font_name='Courier New', space_after=1)

# Cross-feature impact
card2 = add_card(slide, 7.3, 1.6, 5.3, 3.0, WARN)
tb3 = add_textbox(slide, 7.6, 1.8, 4.8, 0.4)
add_text(tb3.text_frame, '⚠️  크로스-피처 영향 감지', size=16, color=WARN, bold=True)
impact = (
    '"login" 수정 시 다음 피처에 영향:\n\n'
    '• cart\n'
    '  login 인증 토큰 의존\n'
    '  → 인증 로직 변경 시 cart도 수정\n\n'
    '• payment\n'
    '  cart 경유로 login에 간접 의존'
)
tb4 = add_textbox(slide, 7.6, 2.3, 4.8, 2.2)
add_text(tb4.text_frame, impact, size=13, color=WHITE, space_after=2)

card3 = add_card(slide, 7.3, 5.0, 5.3, 1.8, ACCENT3)
tb5 = add_textbox(slide, 7.6, 5.1, 4.8, 0.4)
add_text(tb5.text_frame, '안전 장치', size=16, color=ACCENT3, bold=True)
tb6 = add_textbox(slide, 7.6, 5.5, 4.8, 1.2)
add_text(tb6.text_frame, '→ fix → check → fix 재귀 방지\n→ 영향 피처 사용자 확인\n→ 모든 변경 memory에 기록', size=14, color=WHITE, space_after=3)

# ════════════════════════════════════════════
# SLIDE 21: v0.11.0 Changes
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, 'v0.11.0 변경 사항', size=32, color=ACCENT, bold=True)
tb_sub = add_textbox(slide, 0.8, 0.9, 10, 0.5)
add_text(tb_sub.text_frame, '코드 리뷰 반영 + 게이트 핑퐁 루프', size=18, color=MUTED)

# Security table
tbl = add_table(slide, 0.8, 1.6, 5.5, 6, 2, [1.8, 3.7])
for j, h in enumerate(['항목', '변경']):
    set_cell(tbl, 0, j, h, ACCENT, bold=True)
sec_items = [
    ('bash-guard', 'sudo, 분리플래그, $HOME, chmod 777 차단'),
    ('doc-tracker', 'includes() → endsWith() 오탐 수정'),
    ('prompt-handler', '복합 키워드로 false positive 감소'),
    ('status.js', 'updatePhase 재귀 → 재로드 방식'),
    ('에러 로깅', '모든 catch에 debugLog 추가'),
]
for i, (a, b) in enumerate(sec_items):
    set_cell(tbl, i+1, 0, a, ACCENT3)
    set_cell(tbl, i+1, 1, b)

# Feature table
tbl2 = add_table(slide, 6.8, 1.6, 5.8, 7, 2, [2.0, 3.8])
for j, h in enumerate(['항목', '변경']):
    set_cell(tbl2, 0, j, h, ACCENT, bold=True)
feat_items = [
    ('게이트 루프', '핑퐁 피드백 루프 추가'),
    ('no-op 훅', 'notification/subagent-stop 제거'),
    ('타이포그래피', 'design 템플릿 기본값 + H3/Small'),
    ('OWASP', 'review 템플릿 A06/A08/A09/A10 추가'),
    ('버전 통일', 'plugin/marketplace/config → 0.11.0'),
    ('변경 이력', 'AGENTS.md 시맨틱 버전 순 정렬'),
]
for i, (a, b) in enumerate(feat_items):
    set_cell(tbl2, i+1, 0, a, ACCENT2)
    set_cell(tbl2, i+1, 1, b)

# Test badge
card = add_card(slide, 0.8, 5.5, 11.8, 1.0, ACCENT3)
tb3 = add_textbox(slide, 1.1, 5.6, 11.3, 0.8)
add_text(tb3.text_frame, '✅  테스트: 25 suites, 67 tests — 전체 통과', size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 22: Summary
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 0.8, 0.3, 10, 0.7)
add_text(tb.text_frame, '요약', size=36, color=ACCENT, bold=True)

card_l = add_card(slide, 0.8, 1.3, 5.8, 4.5, ACCENT)
tb2 = add_textbox(slide, 1.1, 1.5, 5.3, 0.5)
add_text(tb2.text_frame, '🎯 VAIS Code가 해결하는 것', size=20, color=ACCENT, bold=True)
problems = '→ AI가 기획 없이 코드부터 작성\n→ 설계와 구현의 괴리 (Gap)\n→ 세션 간 맥락 유실\n→ 피처 간 의존성 파악 불가\n→ 수정 시 영향 범위 모름'
tb3 = add_textbox(slide, 1.1, 2.2, 5.3, 3.0)
add_text(tb3.text_frame, problems, size=16, color=WHITE, space_after=6)

card_r = add_card(slide, 7.0, 1.3, 5.8, 4.5, ACCENT3)
tb4 = add_textbox(slide, 7.3, 1.5, 5.3, 0.5)
add_text(tb4.text_frame, '✨ 핵심 역량', size=20, color=ACCENT3, bold=True)
capabilities = '→ 9단계 파이프라인 — 체계적 산출물\n→ 체이닝 문법 — 유연한 워크플로우\n→ 병렬 에이전트 — design, fe+be\n→ Gap 분석 — 90%+ 자동 검증\n→ Manager Memory — 영속 맥락\n→ 핑퐁 게이트 — 사용자 통제 보장'
tb5 = add_textbox(slide, 7.3, 2.2, 5.3, 3.0)
add_text(tb5.text_frame, capabilities, size=16, color=WHITE, space_after=6)

tb6 = add_textbox(slide, 2, 6.0, 9, 0.5)
add_text(tb6.text_frame, '/vais auto {피처명}  으로 시작하세요', size=22, color=MUTED, alignment=PP_ALIGN.CENTER)
tb7 = add_textbox(slide, 2, 6.6, 9, 0.4)
add_text(tb7.text_frame, 'VAIS Code v0.11.0 · MIT License · github.com/ghlee3401/vais-claude-code', size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 23: Q&A
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
tb = add_textbox(slide, 2, 2.0, 9, 2)
add_text(tb.text_frame, 'Q&A', size=100, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(slide, 2, 4.5, 9, 1)
add_text(tb2.text_frame, '감사합니다', size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
tb3 = add_textbox(slide, 2, 5.8, 9, 0.5)
add_text(tb3.text_frame, 'VAIS Code v0.11.0  ·  /vais help 로 시작하세요', size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Save
# ════════════════════════════════════════════
output_path = '/home/user/vais-claude-code/docs/vais-code-v0.11.0-presentation.pptx'
prs.save(output_path)
print(f'✅ Saved: {output_path}')
print(f'   Slides: {len(prs.slides)}')
